# Quantra AI — general seed pitch deck ($2.5M raise, $2–3M range) → Desktop\Quantra AI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

BG = RGBColor(0x0A, 0x0F, 0x1C); PANEL = RGBColor(0x12, 0x1A, 0x2E)
MINT = RGBColor(0x34, 0xD3, 0x99); CYAN = RGBColor(0x22, 0xD3, 0xEE); INDIGO = RGBColor(0x81, 0x8C, 0xF8)
TEXT = RGBColor(0xE7, 0xEC, 0xF5); MUT = RGBColor(0x93, 0xA0, 0xB8); MUT2 = RGBColor(0x6B, 0x78, 0x90)

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = BG; r.line.fill.background()
    r.shadow.inherit = False
    return s

def txt(s, x, y, w, h, lines, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = tb.text_frame; tf.word_wrap = True
    for i, (t, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; r = p.add_run(); r.text = t
        f = r.font; f.size = Pt(size); f.color.rgb = color; f.bold = bold; f.name = "Segoe UI"
    return tb

def head(s, kicker, title):
    txt(s, 0.7, 0.45, 12, 0.5, [(kicker.upper(), 13, MINT, True)])
    txt(s, 0.7, 0.85, 12, 1.0, [(title, 33, TEXT, True)])

def bullets(s, x, y, w, items, size=16, gap_extra=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(5.2)); tf = tb.text_frame; tf.word_wrap = True
    for i, (lead, rest) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        r1 = p.add_run(); r1.text = "•  " + lead; r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = TEXT; r1.font.name = "Segoe UI"
        if rest:
            r2 = p.add_run(); r2.text = " — " + rest; r2.font.size = Pt(size); r2.font.color.rgb = MUT; r2.font.name = "Segoe UI"

def card(s, x, y, w, h, title, value, sub, accent=MINT):
    r = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = PANEL; r.line.color.rgb = accent; r.line.width = Pt(1); r.shadow.inherit = False
    tf = r.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    a = p.add_run(); a.text = value; a.font.size = Pt(26); a.font.bold = True; a.font.color.rgb = accent; a.font.name = "Segoe UI"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    b = p2.add_run(); b.text = title; b.font.size = Pt(13); b.font.bold = True; b.font.color.rgb = TEXT; b.font.name = "Segoe UI"
    if sub:
        p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
        c = p3.add_run(); c.text = sub; c.font.size = Pt(10.5); c.font.color.rgb = MUT2; c.font.name = "Segoe UI"

def foot(s, note="Probabilistic analytics — never guaranteed returns. Not investment advice."):
    txt(s, 0.7, 7.05, 12, 0.4, [(note, 10, MUT2, False)])

# 1 — title
s = slide()
txt(s, 0.7, 2.3, 12, 1.2, [("Quantra AI", 60, TEXT, True)], PP_ALIGN.CENTER)
txt(s, 0.7, 3.5, 12, 0.7, [("The honest markets terminal — live AI analysis, publicly graded against reality", 22, MINT, False)], PP_ALIGN.CENTER)
txt(s, 0.7, 4.5, 12, 0.6, [("Seed round · raising US$ 2.5M ($2–3M range)  |  quantra-ai.onrender.com", 16, MUT, False)], PP_ALIGN.CENTER)
txt(s, 0.7, 5.1, 12, 0.5, [("Eshan Thanvi, Founder · eshanthanvi@gmail.com", 14, MUT2, False)], PP_ALIGN.CENTER)
foot(s)

# 2 — problem
s = slide(); head(s, "The problem", "Retail investors get noise, or lies")
bullets(s, 0.7, 2.1, 11.9, [
    ("A professional terminal costs $25,000+/year", "so the fast-growing retail class in MENA & South Asia runs on fragmented apps and Telegram tipsters."),
    ("The signals industry sells certainty that cannot exist", "“guaranteed calls”, back-filled win rates, cherry-picked screenshots — and regulators are closing in on it."),
    ("No mainstream product proves its accuracy", "every app claims to be right; none publish an auditable record of when they were wrong."),
    ("Result", "first-generation investors — millions of new accounts across the GCC and India (NSE reports 100M+ registered investors) — have nowhere trustworthy to learn and act."),
]); foot(s)

# 3 — solution
s = slide(); head(s, "The solution", "One honest terminal — already live")
bullets(s, 0.7, 2.0, 11.9, [
    ("Everything on one screen", "crypto (tick-by-tick), stocks on 24 exchanges (incl. DFM, Tadawul, NSE), ETFs, commodities, CME futures, indices, FX, US options chains, and a full Web3 on-chain dashboard."),
    ("AI that explains, not commands", "plain-language analysis, a probabilistic Quantra Score, calibrated projection bands, a whole-market movers radar (odds of ±2–40% moves, 1h–30d), opt-in alerts."),
    ("Radical honesty as architecture", "every projection is a probability with the downside shown — and every one is publicly graded on a tamper-evident ledger."),
    ("Shipped", "web + installable PWA + Android; accounts required; plans & billing built; self-testing infrastructure with CI on every release."),
]); foot(s)

# 4 — proof
s = slide(); head(s, "The moat", "Trust you can audit — nobody can fake this quickly")
card(s, 0.7, 2.3, 3.9, 1.9, "days of daily snapshots", "26", "grows every day, hash-chained")
card(s, 4.75, 2.3, 3.9, 1.9, "graded projections", "17,136", "scored vs real outcomes", CYAN)
card(s, 8.8, 2.3, 3.9, 1.9, "calibration observations", "15,800+", "bands self-correct to 80% truth", INDIGO)
bullets(s, 0.7, 4.6, 11.9, [
    ("Tamper-evident", "each day's record is SHA-256 chained to the last — back-dating or editing breaks the chain, publicly."),
    ("Self-calibrating", "measured band coverage feeds back into the engine, so an “80% band” converges to a true 80%."),
    ("Why it's a moat", "a competitor can copy features in months; they cannot copy months of public, verified accuracy history."),
]); foot(s)

# 5 — market
s = slide(); head(s, "The market", "A generational retail-investing wave in our home markets")
bullets(s, 0.7, 2.1, 11.9, [
    ("GCC", "record retail participation on Tadawul and DFM, driven by national digitization programs and IPO pipelines."),
    ("India & South Asia", "NSE registered investors publicly reported above 100M and compounding; F&O volumes are the world's largest by contracts."),
    ("Crypto & Web3", "MENA is among the fastest-growing crypto adoption regions (Chainalysis); UAE has purpose-built regulation (VARA/ADGM)."),
    ("Wedge", "start with the underserved multi-asset retail analyst; expand to teams, developer API, and licensed-data premium tiers."),
    ("Business model", "freemium SaaS — Free / Pro / Ultimate tiers already built with usage limits; billing is one switch away."),
]); foot(s)

# 6 — why now / why us
s = slide(); head(s, "Why now, why us", "The window and the unfair advantages")
bullets(s, 0.7, 2.1, 11.9, [
    ("Regulatory tailwind", "as GCC and Indian regulators crack down on fake-signal apps, the transparency-first product becomes the compliant default."),
    ("AI cost curve", "the product was built AI-assisted by a single founder — the capital efficiency VCs' own AI theses describe, demonstrated."),
    ("Complete platform, day one", "8 asset classes, personalization, teams, alerts, paper trading, PWA + Android — the feature surface of a Series-A company at seed cost."),
    ("Verified engineering", "38-test CI suite, self-diagnostics every 12 minutes, documented runbooks (billing, ops, scale) — diligence-ready."),
]); foot(s)

# 7 — use of funds (the bifurcation)
s = slide(); head(s, "Use of funds", "US$ 2.5M · 24-month plan")
rows = [
    ("35%", "$875K", "Team", "4 engineers + 1 data specialist + 1 growth lead; founder ops", MINT),
    ("30%", "$750K", "Data & infrastructure", "licensed exchange feeds (Gulf real-time, NSE/BSE F&O), premium APIs, production-grade hosting & CDN", CYAN),
    ("20%", "$500K", "Growth & go-to-market", "Arabic/RTL launch, GCC + India marketing, app-store releases, referral & share loops", INDIGO),
    ("10%", "$250K", "Regulatory & legal", "ADGM/UAE footing for a paid research product; compliance counsel", MUT),
    ("5%",  "$125K", "Reserve & operations", "contingency, audit, insurance", MUT2),
]
y = 2.05
for pct, amt, t, d, c in rows:
    bar = s.shapes.add_shape(1, Inches(0.7), Inches(y), Inches(1.5), Inches(0.78))
    bar.fill.solid(); bar.fill.fore_color.rgb = PANEL; bar.line.color.rgb = c; bar.line.width = Pt(1.2); bar.shadow.inherit = False
    tf = bar.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = pct; r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = c; r.font.name = "Segoe UI"
    txt(s, 2.45, y + 0.02, 2.0, 0.7, [(amt, 17, TEXT, True)])
    txt(s, 4.3, y - 0.03, 8.4, 0.9, [(t, 15, TEXT, True), (d, 11.5, MUT, False)])
    y += 0.95
foot(s)

# 8 — what the money builds (bifurcation → roadmap)
s = slide(); head(s, "What each dollar unlocks", "Funding → shipped roadmap (already specified in code)")
bullets(s, 0.7, 2.0, 11.9, [
    ("Months 0–3  ·  Data & compliance", "licensed Gulf + NSE/BSE derivative feeds go live (adapters already built); ADGM regulatory process starts; infra upgraded; billing switched on."),
    ("Months 3–9  ·  Growth engine", "Arabic/RTL product, share-loop analysis cards, educational layer, referral program, Play Store / App Store launches, first paid cohorts."),
    ("Months 9–18  ·  Depth & expansion", "user backtesting playground, portfolio risk analytics, developer API tier, Africa market adapters (JSE live already), team/enterprise plans."),
    ("Months 18–24  ·  Series-A position", "revenue traction + the only publicly-audited accuracy record in the category = the defensible data story for the next round."),
]); foot(s)

# 9 — the ask
s = slide()
txt(s, 0.7, 2.2, 12, 1.0, [("The ask", 20, MINT, True)], PP_ALIGN.CENTER)
txt(s, 0.7, 2.8, 12, 1.2, [("US$ 2.5 million seed", 48, TEXT, True)], PP_ALIGN.CENTER)
txt(s, 0.7, 4.0, 12, 0.8, [("to take a live, complete, provably honest markets platform", 18, MUT, False),
                            ("from product-ready to market-winning across MENA & South Asia", 18, MUT, False)], PP_ALIGN.CENTER)
txt(s, 0.7, 5.3, 12, 0.8, [("Live demo: quantra-ai.onrender.com  ·  Public proof: /track-record.html", 15, CYAN, True),
                            ("Eshan Thanvi · eshanthanvi@gmail.com", 14, MUT2, False)], PP_ALIGN.CENTER)
foot(s)

out_dir = os.path.join(os.path.expanduser("~"), "Desktop", "Quantra AI")
os.makedirs(out_dir, exist_ok=True)
out = os.path.join(out_dir, "Quantra_AI_Seed_Deck.pptx")
prs.save(out)
print("SAVED:", out)
