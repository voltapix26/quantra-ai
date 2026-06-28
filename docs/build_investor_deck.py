# -*- coding: utf-8 -*-
"""Quantra AI — strategic investor pitch deck (800k AED) by Eshan Thanvi."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG, PANEL = RGBColor(0x0B, 0x0E, 0x14), RGBColor(0x12, 0x17, 0x20)
INK, MUTE = RGBColor(0xE9, 0xEF, 0xF6), RGBColor(0x8B, 0x98, 0xA8)
GREEN, RED = RGBColor(0x34, 0xD3, 0x99), RGBColor(0xFB, 0x71, 0x85)
BLUE, GOLD = RGBColor(0x5B, 0x8D, 0xEF), RGBColor(0xF5, 0xC9, 0x66)
LINE = RGBColor(0x24, 0x2C, 0x38)
EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation(); prs.slide_width, prs.slide_height = EMU_W, EMU_H
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = BG; r.line.fill.background(); r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def box(s, x, y, w, h, fill=None, line=None, lw=1.0, rad=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rad else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.space_before = Pt(0)
        for (t, sz, c, b) in para:
            r = p.add_run(); r.text = t; r.font.size = Pt(sz); r.font.color.rgb = c; r.font.bold = b; r.font.name = "Segoe UI"
    return tb


def header(s, kicker, title, color=GREEN):
    box(s, 0.85, 0.62, 0.16, 0.42, fill=color, rad=True)
    text(s, 1.15, 0.5, 11.3, 0.4, [[(kicker, 13, color, True)]])
    text(s, 1.13, 0.92, 11.5, 1.0, [[(title, 29, INK, True)]])
    box(s, 1.15, 1.62, 11.05, 0.0, line=LINE, lw=1.2)


def bullets(s, x, y, w, items, gap=0.62, size=15, dot=GREEN):
    cy = y
    for head, sub in items:
        box(s, x, cy + 0.07, 0.12, 0.12, fill=dot, rad=True)
        runs = [[(head, size, INK, True)]]
        if sub: runs.append([(sub, 12.5, MUTE, False)])
        text(s, x + 0.32, cy - 0.04, w - 0.32, gap, runs, space=2); cy += gap


def stat(s, x, y, w, big, label, color=GREEN):
    box(s, x, y, w, 1.5, fill=PANEL, line=LINE, rad=True)
    text(s, x, y + 0.2, w, 0.8, [[(big, 33, color, True)]], align=PP_ALIGN.CENTER)
    text(s, x, y + 0.95, w, 0.5, [[(label, 12, MUTE, False)]], align=PP_ALIGN.CENTER)


# 1 — TITLE
s = slide(); box(s, 0, 0, 13.333, 0.16, fill=GREEN)
text(s, 0.9, 1.7, 11.5, 1.0, [[("Quantra ", 52, INK, True), ("AI", 52, GREEN, True)]])
text(s, 0.95, 2.85, 11.4, 0.6, [[("Strategic Investment Proposal", 22, INK, True)]])
text(s, 0.95, 3.6, 11.4, 0.6, [[("A live, professional-grade market-analysis platform seeking a strategic partner.", 15, MUTE, False)]])
box(s, 0.95, 4.5, 4.0, 0.95, fill=PANEL, line=GREEN, lw=1.4, rad=True)
text(s, 0.95, 4.66, 4.0, 0.7, [[("Raise: ", 15, MUTE, False), ("AED 800,000", 20, GREEN, True)]], align=PP_ALIGN.CENTER)
text(s, 0.95, 6.4, 11.4, 0.5, [[("Eshan Thanvi", 14, INK, True), ("   ·   eshanthanvi@gmail.com   ·   quantra-ai.onrender.com", 13, MUTE, False)]])

# 2 — THE OPPORTUNITY
s = slide(); header(s, "THE OPPORTUNITY", "An institutional-grade platform — already live", GREEN)
text(s, 1.15, 2.0, 11.0, 1.0, [[("Professional market tools (Bloomberg, Refinitiv) cost $20,000+ per user/year and are "
   "out of reach for retail investors, advisors and small funds. Quantra AI delivers the core of "
   "that capability — real-time data, advanced charting, an AI analyst and calibrated forecasts — "
   "from any browser.", 15, MUTE, False)]], space=8)
stat(s, 1.15, 3.6, 2.6, "Live", "in production today", GREEN)
stat(s, 3.95, 3.6, 2.6, "25", "world exchanges", BLUE)
stat(s, 6.75, 3.6, 2.6, "80%", "calibrated forecasts", GOLD)
stat(s, 9.55, 3.6, 2.6, "1", "solo-built platform", RED)
text(s, 1.15, 5.5, 11.0, 0.8, [[("This is not an idea on paper. The product is built, deployed and verifiable — "
   "which removes the single biggest risk in an early-stage investment: execution.", 14, INK, True)]], space=4)

# 3 — PROBLEM
s = slide(); header(s, "THE PROBLEM", "Serious market analysis is expensive and fragmented", RED)
bullets(s, 1.2, 2.05, 11.0, [
    ("Cost wall", "Bloomberg/Refinitiv terminals run $20k+/user/year — inaccessible to most of the market."),
    ("Fragmentation", "Charts, news, portfolio and forecasts live in disconnected tools."),
    ("Decorative forecasts", "Retail 'projections' are uncalibrated — pretty bands with no statistical meaning."),
    ("Opacity", "Delayed feeds and unclear data quality erode trust."),
], gap=0.92)

# 4 — WHAT'S BUILT (de-risked)
s = slide(); header(s, "WHAT'S ALREADY BUILT", "A complete, deployed product — risk removed", GREEN)
bullets(s, 1.2, 2.0, 5.4, [
    ("Global coverage", "25 exchanges, 32 indices, ETFs, crypto, FX, commodities."),
    ("Pro charting", "16 technical studies on line & candlestick charts."),
    ("Calibrated forecasts", "Monte-Carlo bands proven to ~80% coverage, live."),
], gap=1.0)
bullets(s, 6.9, 2.0, 5.3, [
    ("AI analyst & brief", "Grounded per-asset analysis + daily digest."),
    ("Full workflow", "Portfolio, paper trading, alerts, calendar, community."),
    ("Broker-linked trading", "Users connect their own brokerage; Quantra never holds funds."),
], gap=1.0, dot=BLUE)
box(s, 1.2, 5.2, 11.0, 1.05, fill=PANEL, line=LINE, rad=True)
text(s, 1.5, 5.4, 10.4, 0.8, [[("Multi-tenant SaaS · PWA with push · continuous deployment · super-admin oversight. "
   "Real-time US data live via Finnhub; expanding to all markets.", 13.5, MUTE, False)]])

# 5 — TRACTION / PROOF
s = slide(); header(s, "PROOF", "Verifiable, not promised", BLUE)
stat(s, 1.15, 2.05, 3.4, "80.1%", "realised forecast-band coverage", GREEN)
stat(s, 4.75, 2.05, 3.4, "Real-time", "US equities (Finnhub), exact", BLUE)
stat(s, 8.35, 2.05, 3.4, "10/10", "core roadmap shipped", GOLD)
text(s, 1.15, 3.95, 11.0, 1.8, [
    [("Calibration is published live ", 15, INK, True),
     ("on a public Track Record page — the model's forecast accuracy is measured against real "
      "outcomes, not back-filled.", 14, MUTE, False)],
    [("The platform is open to inspect ", 15, INK, True),
     ("at quantra-ai.onrender.com — an investor can validate every claim before committing.", 14, MUTE, False)],
], space=10)

# 6 — THE OFFER: EXCLUSIVITY
s = slide(); header(s, "THE OFFER TO YOU", "Strategic exclusivity for AED 800,000", GOLD)
text(s, 1.15, 2.0, 11.0, 0.8, [[("In return for the investment, the partner secures an exclusive strategic position in Quantra AI. "
   "Structure is negotiable — designed around your goals:", 15, MUTE, False)]], space=6)
bullets(s, 1.2, 3.0, 11.0, [
    ("Equity stake", "A meaningful ownership share in the company (percentage by negotiation/valuation)."),
    ("Exclusive licence", "White-label / sector or regional exclusivity — Quantra delivered exclusively to your group or clients."),
    ("Board / advisory seat", "Direct influence over roadmap and strategy."),
    ("First right", "Priority on future rounds and on co-developed products."),
], gap=0.82)

# 7 — USE OF FUNDS (bifurcation)
s = slide(); header(s, "USE OF FUNDS", "How the AED 800,000 is deployed", GREEN)
rows = [("Engineering & product", "Broker integrations, automation engine, mobile apps", "AED 240,000", "30%", GREEN),
        ("Regulatory & compliance", "Licensing + legal review across target markets", "AED 180,000", "22.5%", RED),
        ("Market-data licensing", "Real-time aggregators now; enterprise feeds as we scale", "AED 160,000", "20%", BLUE),
        ("Infrastructure & security", "Scaling, data ops, security hardening", "AED 120,000", "15%", GOLD),
        ("Go-to-market & growth", "Onboarding, partnerships, acquisition", "AED 100,000", "12.5%", MUTE)]
cy = 2.0
for t, d, amt, pct, c in rows:
    box(s, 1.15, cy, 11.05, 0.82, fill=PANEL, line=LINE, rad=True)
    box(s, 1.4, cy + 0.22, 0.14, 0.4, fill=c, rad=True)
    text(s, 1.7, cy + 0.13, 6.6, 0.6, [[(t, 14.5, INK, True), ("   " + d, 11, MUTE, False)]])
    text(s, 8.4, cy + 0.18, 2.4, 0.5, [[(amt, 14, INK, True)]], align=PP_ALIGN.RIGHT)
    text(s, 11.0, cy + 0.18, 1.0, 0.5, [[(pct, 14, c if c != MUTE else MUTE, True)]], align=PP_ALIGN.RIGHT)
    cy += 0.92

# 8 — DATA STRATEGY (honest Bloomberg framing)
s = slide(); header(s, "DATA STRATEGY", "From real-time aggregators to enterprise feeds", BLUE)
bullets(s, 1.2, 2.05, 11.0, [
    ("Today", "Finnhub (real-time US), plus aggregator feeds (Twelve Data / Polygon) for global coverage — already wired."),
    ("With funding", "Premium real-time licensing across all target markets, including Gulf & Asia exchanges."),
    ("At scale", "Enterprise data agreements (Bloomberg Data License / Refinitiv) once volume and redistribution licensing justify the cost."),
], gap=0.95)
box(s, 1.2, 5.0, 11.0, 1.15, fill=PANEL, line=GOLD, lw=1.2, rad=True)
text(s, 1.5, 5.2, 10.4, 0.9, [[("Honest note: ", 13.5, GOLD, True),
   ("Bloomberg/Refinitiv are enterprise contracts (six-figure, with redistribution + exchange fees), not a self-serve key. "
    "We adopt them when scale makes them economical — funding accelerates that path.", 13, MUTE, False)]])

# 9 — AUTOMATED EXECUTION + LATENCY ROADMAP (honest tiers, no guarantee)
s = slide(); header(s, "AUTOMATED EXECUTION", "From broker-API speed to an ultra-low-latency tier", GOLD)
text(s, 1.15, 1.95, 11.0, 0.5, [[("Built on the 80%-calibrated models — a staged latency roadmap, with capital unlocking each tier:", 14, MUTE, False)]], space=4)
tiers = [("Today — built", "Broker-linked automated execution; rules-based orders route through the firm's own regulated brokerage (Quantra never holds funds). Latency: sub-second.", GREEN),
         ("Next — co-located engine", "Execution engine co-located near broker/exchange gateways to cut round-trip toward millisecond-class. Funded by this round.", BLUE),
         ("Vision — ultra-low-latency tier", "A dedicated co-located, FPGA-accelerated execution layer targeting microsecond-class order routing — the institutional HFT frontier.", GOLD)]
cy = 2.55
for t, d, c in tiers:
    box(s, 1.15, cy, 11.05, 1.0, fill=PANEL, line=LINE, rad=True)
    box(s, 1.4, cy + 0.28, 0.14, 0.45, fill=c, rad=True)
    text(s, 1.7, cy + 0.16, 10.2, 0.8, [[(t + "  —  ", 14.5, c, True), (d, 12, MUTE, False)]], space=2)
    cy += 1.1
box(s, 1.15, 5.95, 11.05, 0.95, fill=PANEL, line=RED, lw=1.2, rad=True)
text(s, 1.45, 6.12, 10.5, 0.7, [[("Honest scope: ", 12.5, RED, True),
   ("Microsecond-class execution requires exchange co-location and specialised hardware — it is a future, capital-intensive "
    "phase, not what a cloud app does today, and not high-frequency trading at launch. No strategy guarantees profit.", 11.5, MUTE, False)]])

# 10 — MARKET & MODEL
s = slide(); header(s, "MARKET & MODEL", "Large market, recurring revenue", GREEN)
bullets(s, 1.2, 2.05, 5.4, [
    ("Who", "Retail investors, advisors, students and small funds priced out of incumbents."),
    ("Revenue", "Subscription tiers (Free / Pro / Ultimate) via Stripe — already integrated."),
], gap=1.2)
bullets(s, 6.9, 2.05, 5.3, [
    ("Expansion", "Broker-linked trading, white-label, regional editions."),
    ("Moat", "Calibrated (provable) forecasts + data transparency + AI analyst at a fraction of incumbent cost."),
], gap=1.2, dot=BLUE)

# 11 — FOUNDER
s = slide(); header(s, "FOUNDER", "Built end-to-end by one operator", BLUE)
text(s, 1.15, 2.1, 11.0, 2.2, [
    [("Eshan Thanvi", 18, INK, True)],
    [("Designed, built and deployed the entire Quantra AI platform — data engineering, analytics, forecasting, "
      "the full web application, multi-tenant accounts, billing and the broker integration — and took it live in "
      "production.", 15, MUTE, False)],
    [("", 6, MUTE, False)],
    [("The capital accelerates what one person has already proven: it funds the team, data licensing, compliance "
      "and go-to-market to scale a working product.", 15, INK, True)],
], space=10)

# 12 — RISK (honest)
s = slide(); header(s, "RISK & TRANSPARENCY", "An honest view", RED)
bullets(s, 1.2, 2.05, 11.0, [
    ("Markets carry risk", "Quantra provides analysis and tools — not guaranteed returns. No forecast is a promise."),
    ("Regulatory dependency", "Trading features roll out market-by-market, behind the required licences and legal review."),
    ("Data costs", "Real-time depth scales with funding; enterprise feeds are adopted only when economical."),
    ("Why this matters", "We'd rather show you the real risks than overpromise — that is how we intend to run the company."),
], gap=0.92)

# 13 — THE ASK / CLOSE
s = slide(); box(s, 0, 0, 13.333, 0.16, fill=GREEN)
text(s, 0.9, 1.9, 11.5, 1.0, [[("The ask", 40, INK, True)]])
text(s, 0.95, 3.0, 11.4, 1.4, [
    [("AED 800,000", 30, GREEN, True), ("  for a strategic, exclusive position in a live, de-risked platform.", 18, INK, True)],
    [("", 8, MUTE, False)],
    [("Validate every claim at quantra-ai.onrender.com, then let's design the deal around your goals.", 15, MUTE, False)],
], space=8)
box(s, 0.95, 5.2, 4.6, 0.75, fill=GREEN, rad=True)
text(s, 0.95, 5.37, 4.6, 0.4, [[("Eshan Thanvi · eshanthanvi@gmail.com", 13, BG, True)]], align=PP_ALIGN.CENTER)
text(s, 0.97, 6.5, 11.0, 0.5, [[("Confidential — for the named recipient only. Not an offer of securities; subject to definitive agreement.", 10, MUTE, False)]])

prs.save(r"C:\Users\eshan\quantra-terminal\docs\Quantra_AI_Investor_Proposal.pptx")
print("OK", len(prs.slides._sldIdLst), "slides")
