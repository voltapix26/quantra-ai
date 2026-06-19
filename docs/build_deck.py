# -*- coding: utf-8 -*-
"""Builds the Quantra AI pitch deck (docs/Quantra_AI_Pitch_Deck.pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Brand palette ----
BG      = RGBColor(0x0B, 0x0E, 0x14)   # near-black navy
PANEL   = RGBColor(0x12, 0x17, 0x20)   # card
INK     = RGBColor(0xE9, 0xEF, 0xF6)   # near-white
MUTE    = RGBColor(0x8B, 0x98, 0xA8)   # muted grey
GREEN   = RGBColor(0x34, 0xD3, 0x99)   # up / brand
RED     = RGBColor(0xFB, 0x71, 0x85)   # down
BLUE    = RGBColor(0x5B, 0x8D, 0xEF)   # accent
GOLD    = RGBColor(0xF5, 0xC9, 0x66)
LINE    = RGBColor(0x24, 0x2C, 0x38)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = EMU_W, EMU_H
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = BG; r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, radius=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4):
    """runs: list of paragraphs; each paragraph is list of (txt,size,color,bold)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.space_before = Pt(0)
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
            r.font.name = "Segoe UI"
    return tb


def tag(s, x, y, label, color=GREEN):
    t = box(s, x, y, 0.16, 0.42, fill=color, radius=True)
    return t


def header(s, kicker, title, color=GREEN):
    tag(s, 0.85, 0.62, kicker, color)
    text(s, 1.15, 0.5, 11.3, 0.5, [[(kicker, 13, color, True)]])
    text(s, 1.13, 0.92, 11.4, 1.0, [[(title, 30, INK, True)]])
    box(s, 1.15, 1.62, 11.05, 0.0, line=LINE, line_w=1.2)


def bullets(s, x, y, w, items, gap=0.62, size=15, dot=GREEN):
    cy = y
    for head, sub in items:
        box(s, x, cy + 0.07, 0.12, 0.12, fill=dot, radius=True)
        runs = [[(head, size, INK, True)]]
        if sub:
            runs.append([(sub, 12.5, MUTE, False)])
        text(s, x + 0.32, cy - 0.04, w - 0.32, gap, runs, space=2)
        cy += gap


def stat(s, x, y, w, big, label, color=GREEN):
    box(s, x, y, w, 1.55, fill=PANEL, line=LINE, line_w=1.0, radius=True)
    text(s, x, y + 0.22, w, 0.8, [[(big, 40, color, True)]], align=PP_ALIGN.CENTER)
    text(s, x, y + 1.02, w, 0.5, [[(label, 12.5, MUTE, False)]], align=PP_ALIGN.CENTER)


# ============================================================ 1 · TITLE
s = slide()
box(s, 0, 0, 13.333, 0.16, fill=GREEN)
text(s, 0.9, 2.25, 11.5, 1.2, [[("Quantra ", 60, INK, True), ("AI", 60, GREEN, True)]])
text(s, 0.95, 3.45, 11.4, 0.7,
     [[("A professional-grade market-analysis terminal — global markets, real-time data,", 19, MUTE, False)],
      [("calibrated forecasting and an AI analyst, in one multi-tenant web app.", 19, MUTE, False)]])
for i, (lbl, col) in enumerate([("25 world exchanges", GREEN), ("32 global indices", BLUE),
                                ("16-study charting", GOLD), ("80%-calibrated forecasts", RED)]):
    box(s, 0.95 + i * 3.0, 4.75, 2.8, 0.62, fill=PANEL, line=LINE, radius=True)
    text(s, 0.95 + i * 3.0, 4.9, 2.8, 0.4, [[(lbl, 12.5, col, True)]], align=PP_ALIGN.CENTER)
text(s, 0.95, 6.55, 11.4, 0.5,
     [[("quantra-ai.onrender.com", 13, GREEN, True), ("     ·     Build v39  ·  2026", 13, MUTE, False)]])

# ============================================================ 2 · PROBLEM
s = slide()
header(s, "THE PROBLEM", "Serious market analysis is fragmented and expensive", RED)
bullets(s, 1.2, 2.05, 11.0, [
    ("Bloomberg / Refinitiv cost $20k+ a year", "Out of reach for retail traders, students and small funds — and not self-serve."),
    ("Free tools are siloed and shallow", "Charts in one app, news in another, a portfolio in a spreadsheet, forecasts nowhere."),
    ("“Projections” are usually decorative", "Most retail forecasts are uncalibrated — pretty bands with no statistical meaning."),
    ("Data quality is opaque", "Delayed feeds, wrong daily %s and “open/closed” errors erode trust silently."),
], gap=0.92)

# ============================================================ 3 · SOLUTION
s = slide()
header(s, "THE SOLUTION", "One terminal — coverage, intelligence and honesty")
bullets(s, 1.2, 2.05, 5.4, [
    ("Global coverage", "Stocks across 25 exchanges, 32 indices, ETFs, crypto, FX, commodities."),
    ("Real-time where it counts", "Finnhub live US + Coinbase crypto; transparent freshness everywhere."),
    ("Pro charting", "16 studies on line & candlesticks, replay, drawing tools, saved layouts."),
], gap=1.0)
bullets(s, 6.9, 2.05, 5.3, [
    ("Calibrated forecasts", "Monte-Carlo bands proven to ~80% coverage on live data."),
    ("AI analyst + brief", "Grounded chat per asset and a personalized daily digest."),
    ("Full workflow", "Portfolio, paper trading, alerts, calendar, community — account-synced."),
], gap=1.0, dot=BLUE)

# ============================================================ 4 · COVERAGE (stats)
s = slide()
header(s, "MARKET COVERAGE", "From Wall Street to Dalal Street — in native currency", BLUE)
stat(s, 1.15, 2.1, 2.6, "25", "world stock exchanges", GREEN)
stat(s, 3.95, 2.1, 2.6, "32", "global indices", BLUE)
stat(s, 6.75, 2.1, 2.6, "40+", "regional & sector ETFs", GOLD)
stat(s, 9.55, 2.1, 2.6, "16", "display currencies", RED)
text(s, 1.15, 4.05, 11.0, 1.8, [
    [("Bifurcated by exchange — ", 15, INK, True),
     ("Reliance routes to BSE/NSE in ₹, Shell to LSE in £, Emaar to DFM in AED; the display "
      "currency follows the exchange automatically, with pence/cents minor-unit handling.", 14, MUTE, False)],
    [("Asset classes — ", 15, INK, True),
     ("equities, indices, ETFs, cryptocurrencies, FX pairs and commodities, each with the right "
      "session model (24/7 crypto, 24/5 FX, exchange hours for equities).", 14, MUTE, False)],
], space=10)

# ============================================================ 5 · CHARTING
s = slide()
header(s, "CHARTING & INDICATORS", "A 16-study technical suite on every chart type", GOLD)
bullets(s, 1.2, 2.0, 5.4, [
    ("Overlays", "Bollinger, Ichimoku Cloud, Supertrend, VWAP, Keltner, Donchian, EMA, McGinley, Parabolic SAR."),
    ("Oscillator panes", "RSI, MACD, Stochastic, ADX, Williams %R, CCI, Volume."),
], gap=1.25)
bullets(s, 6.9, 2.0, 5.3, [
    ("Chart types", "Line, area, candlesticks, Heikin Ashi — all studies render on each."),
    ("Pro tools", "Bar replay, trendline / level / Fibonacci drawing, enlarge mode, saved layouts."),
], gap=1.25, dot=GOLD)
box(s, 1.2, 4.7, 11.0, 1.4, fill=PANEL, line=LINE, radius=True)
text(s, 1.5, 4.9, 10.4, 1.1, [
    [("Shared overlay engine ", 14, GOLD, True),
     ("— one renderer draws every study identically on line and candlestick charts, with a "
      "dynamic legend that lists only the active overlays and their colours. Studies persist "
      "locally and travel inside account-synced saved layouts.", 13.5, MUTE, False)]])

# ============================================================ 6 · FORECASTING (differentiator)
s = slide()
header(s, "DIFFERENTIATOR", "Forecasts that are actually calibrated", RED)
stat(s, 1.15, 2.05, 3.4, "80.1%", "realised band coverage", GREEN)
stat(s, 4.75, 2.05, 3.4, "74.4%", "before recalibration", MUTE)
stat(s, 8.35, 2.05, 3.4, "Live", "calibration widget", BLUE)
text(s, 1.15, 3.95, 11.0, 2.0, [
    [("Method — ", 15, INK, True),
     ("bootstrap Monte-Carlo from real returns (fat tails), drift shrunk toward a random walk, "
      "and an MC-median central path. The P10–P90 band targets 80% and hits it.", 14, MUTE, False)],
    [("Provable, not back-filled — ", 15, INK, True),
     ("daily snapshots store sigma; the server measures realised in-band coverage from real forward "
      "outcomes and the Track Record page shows live coverage vs the 80% target.", 14, MUTE, False)],
], space=10)

# ============================================================ 7 · AI
s = slide()
header(s, "INTELLIGENCE", "An AI analyst on every asset", BLUE)
bullets(s, 1.2, 2.05, 11.0, [
    ("Ask Quantra", "Grounded conversational analyst per asset — LLM when keyed, a smart local fallback otherwise, with suggested prompts and rate limiting."),
    ("AI Daily Brief", "Personalized digest across your watchlist & portfolio: movers, market breadth and upcoming earnings, narrated."),
    ("Smart alerts", "Server-side engine fires on price/percent conditions even when the tab is closed, delivered by email and web-push."),
], gap=1.05)

# ============================================================ 8 · DATA & TRUST
s = slide()
header(s, "DATA & TRUST", "Honesty as a product feature", GREEN)
bullets(s, 1.2, 2.0, 5.5, [
    ("Real-time where it matters", "Finnhub live US quotes, Coinbase crypto WS → SSE relay (keys stay server-side)."),
    ("Accurate market status", "Open/closed from session time-of-day in each exchange timezone + a per-exchange holiday calendar."),
], gap=1.25)
bullets(s, 6.95, 2.0, 5.3, [
    ("Transparent freshness", "“Live / delayed feed / at last close” labels so a delayed feed never looks like a bug."),
    ("Correct by construction", "Daily %s use prior-session close; absolute change shown beside the %; audited end-to-end."),
], gap=1.25, dot=GREEN)
box(s, 1.2, 4.75, 11.0, 1.3, fill=PANEL, line=LINE, radius=True)
text(s, 1.5, 4.95, 10.4, 1.0, [
    [("Every paid source is key-gated with graceful fallback ", 14, GREEN, True),
     ("— the app degrades to free feeds instead of breaking, and clearly shows when data is delayed.", 13.5, MUTE, False)]])

# ============================================================ 9 · WORKFLOW
s = slide()
header(s, "FULL WORKFLOW", "Research → decide → track — in one place", GOLD)
cards = [("Portfolio", "Holdings, live P&L, allocation & totals", GREEN),
         ("Paper trading", "Simulated cash, positions, realized P&L, journal", BLUE),
         ("Discovery", "Heatmap, movers & breadth per exchange", GOLD),
         ("Calendar", "Earnings, IPOs & economic events", RED),
         ("Community", "Shared ideas, upvotes & leaderboard", GREEN),
         ("PWA + push", "Installable app, offline shell, notifications", BLUE)]
for i, (t, d, c) in enumerate(cards):
    x = 1.15 + (i % 3) * 3.75; y = 2.1 + (i // 3) * 1.85
    box(s, x, y, 3.5, 1.65, fill=PANEL, line=LINE, radius=True)
    box(s, x + 0.28, y + 0.32, 0.14, 0.5, fill=c, radius=True)
    text(s, x + 0.55, y + 0.28, 2.8, 0.5, [[(t, 16, INK, True)]])
    text(s, x + 0.3, y + 0.85, 3.0, 0.7, [[(d, 12.5, MUTE, False)]])

# ============================================================ 10 · ARCHITECTURE
s = slide()
header(s, "ARCHITECTURE", "Lean, resilient, multi-tenant", BLUE)
bullets(s, 1.2, 2.05, 11.0, [
    ("Single Node.js server", "Acts as market-data proxy + app server; optional deps (pg, stripe, web-push) load only when configured."),
    ("Multi-tenant accounts", "Per-user watchlists, portfolios, alerts, layouts & paper accounts; Postgres or file store; super-admin oversight (metadata only — never passwords)."),
    ("Shared analysis engine", "One analytics core (window.Quantra) runs client-side and is reused server-side via a shim."),
    ("Always-on", "Keep-warm self-ping removes free-tier cold starts; SVG charts render crisp at any size."),
], gap=0.92)

# ============================================================ 11 · STATUS / TRACTION
s = slide()
header(s, "STATUS", "Shipped, deployed, and verified", GREEN)
stat(s, 1.15, 2.05, 2.7, "10/10", "roadmap features live", GREEN)
stat(s, 4.05, 2.05, 2.7, "v39", "current build", BLUE)
stat(s, 6.95, 2.05, 2.7, "50+", "deploys this cycle", GOLD)
stat(s, 9.85, 2.05, 2.7, "1", "unified terminal", RED)
text(s, 1.15, 3.95, 11.0, 1.8, [
    [("Live in production ", 15, INK, True),
     ("on Render with GitHub auto-deploy. The full 10-feature SaaS roadmap is complete, followed by "
      "global-market expansion, the indicator suite, real-time feeds and holiday-aware market status.", 14, MUTE, False)],
    [("See CHANGELOG.md ", 15, GREEN, True),
     ("for the complete version history (Phases 1–6, v1.0.0 → v5.7.0).", 14, MUTE, False)],
], space=10)

# ============================================================ 12 · CLOSING
s = slide()
box(s, 0, 0, 13.333, 0.16, fill=GREEN)
text(s, 0.9, 2.5, 11.5, 1.0, [[("The Bloomberg experience,", 38, INK, True)],
                              [("re-imagined for everyone.", 38, GREEN, True)]])
text(s, 0.95, 4.25, 11.4, 0.6,
     [[("Global markets · real-time data · calibrated forecasts · an AI analyst — one terminal.", 16, MUTE, False)]])
box(s, 0.95, 5.3, 4.2, 0.7, fill=GREEN, radius=True)
text(s, 0.95, 5.46, 4.2, 0.4, [[("quantra-ai.onrender.com", 14, BG, True)]], align=PP_ALIGN.CENTER)
text(s, 5.4, 5.46, 6.0, 0.4, [[("github.com/voltapix26/quantra-ai", 13, MUTE, False)]])

prs.save(r"C:\Users\eshan\quantra-terminal\docs\Quantra_AI_Pitch_Deck.pptx")
print("OK", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
